from pptx import Presentation

def open_and_read():
    prs = Presentation("pyppy.pptx")
    slides = prs.slides

    slides_count = len(slides)
    print(slides_count)

    return slides

def override_content(slide, new_content):
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = new_content

    return slide

