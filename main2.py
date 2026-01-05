from datetime import date
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE

text_top = [
    "Cepat",
    "Stabil",
    "Scalable"
]

text_bottom = [
    "Support Python",
    "Auto Generate PPT",
    "Ready for Report"
]

today = date.today()

def set_bullet_list(shape, items, font_size=14):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for i, text in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = text
        p.level = 0
        p.font.size = Pt(font_size)

def read_and_write_ppt():
    try:
        prs = Presentation("pyppt.pptx")
        slide = prs.slides[1]

        # title
        slide.shapes.title.text = f"Update Kondisi Siklon Tropis. Selasa, {today}"

        # === WAJIB pastikan index placeholder ===
        # debug dulu:
        for p in slide.placeholders:
            print(p.placeholder_format.idx, p.name)

        title = None
        title_description = None
        picture = None
        container1 = None
        container2 = None

        for shape in slide.shapes:

            if shape.has_text_frame and shape.text == 'container1':
                container1 = shape

            if shape.has_text_frame and shape.text == 'container2':
                container2 = shape

            if shape.has_text_frame and shape.text == 'judul':
                shape.text = 'Update Kondisi Siklon Tropis dan Bibit Siklon Tropis'
                title = shape

            if shape.has_text_frame and shape.text == 'tanggal':
                shape.text = f"{today}"
                title_description = shape

            if not shape.has_text_frame and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture = shape

            if shape.has_text_frame:
                print(
                    shape.name,
                    "| placeholder:",
                    shape.is_placeholder,
                    "| has_text:",
                    shape.has_text_frame,
                    "| text:",
                    shape.text
                )
            else:
                print(
                    shape.name,
                    "| placeholder:",
                    shape.is_placeholder,
                    "| has_text:",
                    shape.has_text_frame
                )

        set_bullet_list(container1, text_top)
        set_bullet_list(container2, text_bottom)

        if picture:
            left, top, width, height = picture.left, picture.top, picture.width, picture.height
            slide.shapes._spTree.remove(shape._element)
            slide.shapes.add_picture("new_image.png", left, top, width=width, height=height)

        prs.save("output.pptx")
    except Exception as e:
        print(e)

if __name__ == "__main__":
    read_and_write_ppt()
