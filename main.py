from pptx import Presentation
from pptx.util import Inches

def open_and_read():
    try:
        prs = Presentation("pyppt.pptx")
        print(len(prs.slides))

        # get slide ke 2 index 1
        slide = prs.slides[1]

        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text = "Konten baru ges"

        slide.shapes.title.text = "Judul Baru"
        slide.placeholders[1].text = "Isi baru slide 1"

        # cari image lama
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide.shapes._spTree.remove(shape._element)

                slide.shapes.add_picture(
                    "new_image.png",
                    left, top, width=width, height=height
                )
                break

        prs.save("output.pptx")
    except Exception as e :
        print(e)
        return "Something went wrong"


if __name__ == "__main__":
    open_and_read()